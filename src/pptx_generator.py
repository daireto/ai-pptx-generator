"""PowerPoint presentation generator from JSON."""

import random
import re
from io import BytesIO
from typing import IO, Any

import httpx
import orjson
from PIL import Image
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.shapes.base import BaseShape
from pptx.shapes.graphfrm import GraphicFrame
from pptx.slide import Slide
from pptx.text.text import TextFrame
from pptx.util import Inches, Pt
from starlette.datastructures import Secret

from src.definitions import PIXABAY_API_URL


class PPTXGenerator:
    """Generate a PowerPoint presentation from a JSON string."""

    def __init__(self, pixabay_api_key: Secret | None = None) -> None:
        """Initialize the generator.

        Parameters
        ----------
        pixabay_api_key : Secret | None, optional
            API key for Pixabay, by default None.

        """
        self.prs = Presentation()
        self.slide_width = self.prs.slide_width or Inches(10)
        self.slide_height = self.prs.slide_height or Inches(7.5)
        self.__pixabay_api_key = pixabay_api_key

    def generate(self, json_str: str, output_file: str | IO[bytes]) -> None:
        """Generate a PowerPoint presentation from a JSON string.

        Parameters
        ----------
        json_str : str
            JSON string to generate the presentation from.
        output_file : str | IO[bytes]
            Path to the output file or file-like object.

        """
        data = orjson.loads(json_str)
        self._process(data)
        self.prs.save(output_file)

    def generate_from_file(self, json_path: str, output_file: str | IO[bytes]) -> None:
        """Generate a PowerPoint presentation from a JSON file.

        Parameters
        ----------
        json_path : str
            Path to the JSON file.
        output_file : str | IO[bytes]
            Path to the output file or file-like object.

        """
        with open(json_path, encoding='utf-8') as f:
            data = orjson.loads(f.read())

        self._process(data)
        self.prs.save(output_file)

    def test_json(self, json_str: str) -> None:
        """Test if a JSON string is valid.

        Parameters
        ----------
        json_str : str
            JSON string to test.

        """
        data = orjson.loads(json_str)
        self._process(data)
        self.prs = Presentation()
        self.slide_width = self.prs.slide_width or Inches(10)
        self.slide_height = self.prs.slide_height or Inches(7.5)

    def rgb(self, hex_color: str) -> RGBColor:
        """Convert a hex color to a RGBColor object.

        Parameters
        ----------
        hex_color : str
            Hex color to convert.

        Returns
        -------
        RGBColor
            RGBColor object.

        """
        hex_color = hex_color.lstrip('#')
        return RGBColor(
            int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
        )

    def _apply_background(
        self,
        slide: Slide,
        background: dict[str, Any] | None = None,
        background_color: str | None = None,
    ) -> None:
        fill = slide.background.fill
        if background:
            bkg_type = background.get('type', 'solid')
            if bkg_type == 'solid':
                fill.solid()
                fill.fore_color.rgb = self.rgb(background.get('color', '#FFFFFF'))
            elif bkg_type == 'gradient':
                # Limpiar contenido XML actual
                x_pr = fill._xPr  # noqa: SLF001
                for child in list(x_pr):  # type: ignore
                    x_pr.remove(child)

                grad_fill = OxmlElement('a:gradFill')
                grad_fill.set('flip', 'none')

                lin = OxmlElement('a:lin')
                lin.set('ang', '5400000')  # vertical
                lin.set('scaled', '1')
                grad_fill.append(lin)

                gs_lst = OxmlElement('a:gsLst')

                from_color_hex = background.get('from', 'FFFFFF').lstrip('#')
                to_color_hex = background.get('to', '000000').lstrip('#')

                gs1 = OxmlElement('a:gs')
                gs1.set('pos', '0')
                srgb_clr1 = OxmlElement('a:srgbClr')
                srgb_clr1.set('val', from_color_hex)
                gs1.append(srgb_clr1)

                gs2 = OxmlElement('a:gs')
                gs2.set('pos', '100000')
                srgb_clr2 = OxmlElement('a:srgbClr')
                srgb_clr2.set('val', to_color_hex)
                gs2.append(srgb_clr2)

                gs_lst.append(gs1)
                gs_lst.append(gs2)
                grad_fill.append(gs_lst)

                x_pr.append(grad_fill)
            else:
                fill.solid()
                fill.fore_color.rgb = self.rgb('#FFFFFF')
        elif background_color:
            fill.solid()
            fill.fore_color.rgb = self.rgb(background_color)
        else:
            fill.solid()
            fill.fore_color.rgb = self.rgb('#FFFFFF')

    def _apply_textbox_style(
        self, text_frame: TextFrame, textbox_style: dict[str, Any] | None = None
    ) -> None:
        if not textbox_style:
            return
        if 'margin_top' in textbox_style:
            text_frame.margin_top = Pt(textbox_style['margin_top'])
        if 'margin_bottom' in textbox_style:
            text_frame.margin_bottom = Pt(textbox_style['margin_bottom'])
        if 'margin_left' in textbox_style:
            text_frame.margin_left = Pt(textbox_style['margin_left'])
        if 'margin_right' in textbox_style:
            text_frame.margin_right = Pt(textbox_style['margin_right'])
        # padding no es soportado directamente en pptx

    def _apply_shadow_and_border(
        self, shape: BaseShape, effects: dict[str, Any] | None = None
    ) -> None:
        if not effects:
            return

        shadow = effects.get('shadow')
        border = effects.get('border')

        sp = shape._element  # noqa: SLF001
        sp_pr = sp.find(qn('p:spPr'), namespaces=None)
        if sp_pr is None:
            sp_pr = OxmlElement('p:spPr')
            sp.insert(0, sp_pr)

        if shadow:
            for child in sp_pr.findall(qn('a:effectLst'), namespaces=None):
                sp_pr.remove(child)

            effect_lst = OxmlElement('a:effectLst')
            outer_shdw = OxmlElement('a:outerShdw')
            outer_shdw.set('blurRad', str(int(shadow.get('blur', 5) * 12700)))
            outer_shdw.set(
                'dist',
                str(
                    int(
                        (
                            (
                                shadow.get('offset_x', 0) ** 2
                                + shadow.get('offset_y', 0) ** 2
                            )
                            ** 0.5
                        )
                        * 12700
                    )
                ),
            )
            outer_shdw.set('dir', '5400000')
            outer_shdw.set('rotWithShape', '0')

            srgb_clr = OxmlElement('a:srgbClr')
            srgb_clr.set('val', shadow.get('color', '000000').lstrip('#'))
            alpha = OxmlElement('a:alpha')
            alpha.set('val', str(int(shadow.get('opacity', 0.5) * 100000)))
            srgb_clr.append(alpha)
            outer_shdw.append(srgb_clr)
            effect_lst.append(outer_shdw)
            sp_pr.append(effect_lst)

        # Solo aplicar .line si el objeto lo permite
        if border and not isinstance(shape, GraphicFrame):
            shape.line.color.rgb = self.rgb(border.get('color', '#000000'))  # type: ignore
            shape.line.width = Pt(border.get('width', 1))  # type: ignore

    def _add_textbox(
        self,
        slide: Slide,
        text: str,
        style: dict[str, Any],
        textbox_style: dict[str, Any] | None = None,
        effects: dict[str, Any] | None = None,
        position: dict[str, Any] | None = None,
    ) -> None:
        left = Inches(position.get('left', 0.5)) if position else Inches(0.5)
        top = Inches(position.get('top', 0.5)) if position else Inches(0.5)
        width = Inches(position.get('width', 9)) if position else Inches(9)
        height = Inches(position.get('height', 1)) if position else Inches(1)

        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        self._apply_textbox_style(tf, textbox_style)

        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text

        font = run.font
        font.name = style.get('font_name', 'Calibri')
        font.size = Pt(style.get('font_size', 20))
        font.bold = style.get('bold', False)
        font.italic = style.get('italic', False)
        font.underline = style.get('underline', False)
        if 'color' in style:
            font.color.rgb = self.rgb(style['color'])

        alignment = style.get('alignment', 'left').lower()
        p.alignment = {'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}.get(
            alignment, PP_ALIGN.LEFT
        )

        # Interlineado si está definido
        if 'line_spacing' in style:
            p.line_spacing = style['line_spacing']

        self._apply_shadow_and_border(box, effects)

    def _download_image_from_pixabay(self, query: str) -> Image.Image | None:
        if not self.__pixabay_api_key:
            return None

        query = re.sub(r'[-_]', ' ', query)
        response = httpx.get(
            PIXABAY_API_URL,
            params={'key': str(self.__pixabay_api_key), 'q': query},
            timeout=10,
            verify=False
        )
        if not response.is_success:
            return None

        data = response.json()
        if not data['hits']:
            return None

        image_url = random.choice(data['hits'])['webformatURL']
        image_response = httpx.get(image_url, timeout=10, verify=False)
        if not image_response.is_success:
            return None

        return Image.open(BytesIO(image_response.content))

    def _crop_center(self, img: Image.Image, target_aspect: float) -> Image.Image:
        img_aspect = img.width / img.height
        if img_aspect > target_aspect:
            new_width = int(img.height * target_aspect)
            left = (img.width - new_width) // 2
            right = left + new_width
            box = (left, 0, right, img.height)
        else:
            new_height = int(img.width / target_aspect)
            top = (img.height - new_height) // 2
            bottom = top + new_height
            box = (0, top, img.width, bottom)
        return img.crop(box)

    def _add_image(
        self,
        slide: Slide,
        pixabay_query: str,
        style: dict[str, Any],
        effects: dict[str, Any] | None = None,
        position: dict[str, Any] | None = None,
    ) -> None:
        img = self._download_image_from_pixabay(pixabay_query)
        if not img:
            return

        cover = style.get('cover', False)
        pos = position or {}
        left = Inches(pos.get('left', 0))
        top = Inches(pos.get('top', 0))
        width_in = pos.get('width')
        height_in = pos.get('height')

        if cover and width_in and height_in:
            box_width = Inches(width_in)
            box_height = Inches(height_in)
            target_aspect = box_width / box_height
            cropped_img = self._crop_center(img, target_aspect)

            image_stream = BytesIO()
            cropped_img.save(image_stream, format='PNG')
            image_stream.seek(0)

            pic = slide.shapes.add_picture(
                image_stream, left, top, width=box_width, height=box_height
            )
        else:
            img_width_in = width_in if width_in else self.slide_width.inches / 2
            img_height_in = (
                height_in if height_in else img.height * img_width_in / img.width
            )
            if img_height_in > self.slide_height.inches:
                img_height_in = self.slide_height.inches
                img_width_in = img.width * img_height_in / img.height

            width = Inches(img_width_in)
            height = Inches(img_height_in)

            # Posición por defecto center si no se pasa posición izquierda o derecha
            if pos.get('position') == 'center' or (
                'left' not in pos and 'top' not in pos
            ):
                left = (self.slide_width - width) / 2
                top = (self.slide_height - height) / 2

            image_stream = BytesIO()
            img.save(image_stream, format='PNG')
            image_stream.seek(0)
            pic = slide.shapes.add_picture(
                image_stream, Inches(left), Inches(top), width=width, height=height
            )

        self._apply_shadow_and_border(pic, effects)

    def _add_table(
        self,
        slide: Slide,
        table_data: dict[str, Any],
        position: dict[str, Any] | None = None,
    ) -> None:
        pos = position or {}
        left = Inches(pos.get('left', 0.5))
        top = Inches(pos.get('top', 2.5))
        width = Inches(pos.get('width', 9))
        height = Inches(pos.get('height', 3))

        headers = table_data.get('headers', [])
        rows_data = table_data.get('rows', [])
        rows = len(rows_data) + 1
        cols = len(headers)

        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table

        style = table_data.get('style', {})
        header_fill = style.get('header_fill')
        header_font_color = style.get('header_font_color')
        cell_fill = style.get('cell_fill')
        effects = style.get('effects')

        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            if header_fill:
                cell.fill.solid()
                cell.fill.fore_color.rgb = self.rgb(header_fill)
            if header_font_color:
                for p in cell.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.color.rgb = self.rgb(header_font_color)

        for row_idx, row in enumerate(rows_data):
            for col_idx, val in enumerate(row):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(val)
                if cell_fill:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.rgb(cell_fill)

        # Aplica sombra y borde al contenedor de la tabla (shape)
        if effects:
            self._apply_shadow_and_border(table_shape, effects)

    def _add_chart(
        self,
        slide: Slide,
        chart_data_json: dict[str, Any],
        position: dict[str, Any] | None = None,
    ) -> None:
        chart_type_str = chart_data_json.get('type', 'COLUMN_CLUSTERED').upper()
        chart_type = getattr(
            XL_CHART_TYPE,
            chart_type_str,
            XL_CHART_TYPE.COLUMN_CLUSTERED,
        )

        chart_data = CategoryChartData()
        chart_data.categories = chart_data_json.get('categories', [])

        for serie in chart_data_json.get('series', []):
            chart_data.add_series(serie.get('name', ''), serie.get('values', []))

        pos = position or {}

        MIN_WIDTH = Inches(8)
        MIN_HEIGHT = Inches(4)

        raw_left = pos.get('left', 1)
        raw_top = pos.get('top', 2)
        raw_width = pos.get('width', 8)
        raw_height = pos.get('height', 4.5)

        requested_width = Inches(raw_width)
        requested_height = Inches(raw_height)

        # Dimensiones finales con mínimos
        cx = max(requested_width, MIN_WIDTH)
        cy = max(requested_height, MIN_HEIGHT)

        # Compensar si se aumentó el ancho: centrar respecto al valor original
        if cx > requested_width:
            extra_width = cx - requested_width
            x = Inches(raw_left) - int((extra_width / 2))
        else:
            x = Inches(raw_left)

        y = Inches(raw_top)

        chart_shape = slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data)  # type: ignore
        chart = chart_shape.chart  # type: ignore

        # Título del gráfico
        if chart_data_json.get('title'):
            chart.has_title = True
            chart.chart_title.text_frame.text = chart_data_json['title']
        else:
            chart.has_title = False

        # Etiquetas de datos
        if chart_data_json.get('data_labels', False):
            chart.plots[0].has_data_labels = True

        # Leyenda
        legend_position = chart_data_json.get('legend_position', 'right').lower()
        chart.has_legend = True
        position_map = {
            'top': XL_LEGEND_POSITION.TOP,
            'bottom': XL_LEGEND_POSITION.BOTTOM,
            'left': XL_LEGEND_POSITION.LEFT,
            'right': XL_LEGEND_POSITION.RIGHT,
            'corner': XL_LEGEND_POSITION.CORNER,
        }
        chart.legend.position = position_map.get(
            legend_position, XL_LEGEND_POSITION.RIGHT
        )

    def _process(self, data: dict[str, Any]) -> None:
        title_slide_data = data.get('title_slide', {})
        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        self._apply_background(
            slide,
            title_slide_data.get('background'),
            title_slide_data.get('background_color'),
        )

        for shape, content in zip([title, subtitle], ['title', 'subtitle']):
            style = title_slide_data.get('style', {})
            p = shape.text_frame.paragraphs[0]  # type: ignore
            p.clear()
            run = p.add_run()
            run.text = title_slide_data.get(content, '')
            font = run.font
            font.name = style.get('font_name', 'Calibri')
            font.size = Pt(style.get('font_size', 32))
            font.bold = style.get('bold', False)
            font.italic = style.get('italic', False)
            font.underline = style.get('underline', False)
            if 'color' in style:
                font.color.rgb = self.rgb(style['color'])
            alignment = style.get('alignment', 'center').lower()
            p.alignment = {'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}.get(
                alignment, PP_ALIGN.LEFT
            )

        for slide_data in data.get('slides', []):
            layout = slide_data.get('layout', '')
            elements = slide_data.get('elements', {})
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

            self._apply_background(
                slide, slide_data.get('background'), slide_data.get('background_color')
            )

            if layout == 'text_left_image_right':
                if 'title' in elements:
                    self._add_textbox(
                        slide,
                        elements['title']['text'],
                        elements['title'].get('style', {}),
                        textbox_style=elements['title'].get('textbox_style'),
                        effects=elements['title'].get('effects'),
                        position=elements['title'].get('position'),
                    )
                for para in elements.get('paragraphs', []):
                    self._add_textbox(
                        slide,
                        para['text'],
                        para.get('style', {}),
                        textbox_style=para.get('textbox_style'),
                        effects=para.get('effects'),
                        position=para.get('position'),
                    )
                if 'image' in elements:
                    self._add_image(
                        slide,
                        elements['image']['pixabay_query'],
                        elements['image'].get('style', {}),
                        effects=elements['image'].get('style', {}).get('effects'),
                        position=elements['image'].get('position'),
                    )

            elif layout == 'two_column_text':
                if 'title' in elements:
                    self._add_textbox(
                        slide,
                        elements['title']['text'],
                        elements['title'].get('style', {}),
                        textbox_style=elements['title'].get('textbox_style'),
                        effects=elements['title'].get('effects'),
                        position=elements['title'].get('position'),
                    )
                for col in elements.get('columns', []):
                    for para in col.get('paragraphs', []):
                        self._add_textbox(
                            slide,
                            para['text'],
                            para.get('style', {}),
                            textbox_style=para.get('textbox_style'),
                            effects=para.get('effects'),
                            position=para.get('position'),
                        )

            elif layout == 'image_full_width_with_caption':
                if 'image' in elements:
                    self._add_image(
                        slide,
                        elements['image']['pixabay_query'],
                        elements['image'].get('style', {}),
                        effects=elements['image'].get('style', {}).get('effects'),
                        position=elements['image'].get('position'),
                    )
                if 'caption' in elements:
                    self._add_textbox(
                        slide,
                        elements['caption']['text'],
                        elements['caption'].get('style', {}),
                        textbox_style=elements['caption'].get('textbox_style'),
                        effects=elements['caption'].get('effects'),
                        position=elements['caption'].get('position'),
                    )

            elif layout == 'text_top_table_bottom':
                if 'title' in elements:
                    self._add_textbox(
                        slide,
                        elements['title']['text'],
                        elements['title'].get('style', {}),
                        textbox_style=elements['title'].get('textbox_style'),
                        effects=elements['title'].get('effects'),
                        position=elements['title'].get('position'),
                    )
                for para in elements.get('paragraphs', []):
                    self._add_textbox(
                        slide,
                        para['text'],
                        para.get('style', {}),
                        textbox_style=para.get('textbox_style'),
                        effects=para.get('effects'),
                        position=para.get('position'),
                    )
                if 'table' in elements:
                    self._add_table(
                        slide,
                        elements['table'],
                        position=elements['table'].get('position'),
                    )

            elif layout == 'text_with_bullet_points':
                if 'title' in elements:
                    self._add_textbox(
                        slide,
                        elements['title']['text'],
                        elements['title'].get('style', {}),
                        textbox_style=elements['title'].get('textbox_style'),
                        effects=elements['title'].get('effects'),
                        position=elements['title'].get('position'),
                    )
                for bullet in elements.get('bullet_points', []):
                    self._add_textbox(
                        slide,
                        '• ' + bullet['text'],
                        bullet.get('style', {}),
                        textbox_style=bullet.get('textbox_style'),
                        effects=bullet.get('effects'),
                        position=bullet.get('position'),
                    )
            elif layout == 'chart':
                if 'title' in elements:
                    self._add_textbox(
                        slide,
                        elements['title']['text'],
                        elements['title'].get('style', {}),
                        textbox_style=elements['title'].get('textbox_style'),
                        effects=elements['title'].get('effects'),
                        position=elements['title'].get('position'),
                    )
                if 'chart' in elements:
                    self._add_chart(
                        slide, elements['chart'], elements['chart'].get('position')
                    )
